from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


NAVY = RGBColor(24, 50, 74)
BLUE = RGBColor(40, 120, 181)
PALE = RGBColor(234, 242, 248)
GREY = RGBColor(243, 245, 247)
MID = RGBColor(101, 113, 124)
INK = RGBColor(29, 37, 44)
WHITE = RGBColor(255, 255, 255)
AMBER = RGBColor(196, 124, 0)
RED = RGBColor(163, 58, 50)
GREEN = RGBColor(39, 103, 73)
LINE = RGBColor(216, 222, 227)

GUIDANCE = {
    "Paint by Numbers": {
        "plain": "Outcome and method are known. Control the sequence and manage variance.",
        "delivery": ["Baseline", "Standardise", "Check quality", "Control exceptions"],
        "control": ["Milestones", "Acceptance checks", "Dependency log", "Exception report"],
    },
    "Making Movies": {
        "plain": "The result is clear. Quality depends on specialists, integration and review.",
        "delivery": ["Protect vision", "Plan hand-offs", "Prototype", "Review integration"],
        "control": ["Clear brief", "Creative authority", "Stage reviews", "Approvals"],
    },
    "Quest": {
        "plain": "The destination is clear. Test routes before committing to the full plan.",
        "delivery": ["Frame outcome", "Test options", "Compare evidence", "Choose route"],
        "control": ["Hypotheses", "Evidence reviews", "Option comparison", "Decision gates"],
    },
    "Walking in Fog": {
        "plain": "The destination and route are not firm. Create clarity before locking the plan.",
        "delivery": ["Short horizons", "Visible assumptions", "Evidence gates", "Protect options"],
        "control": ["Outcome frame", "Assumption log", "Decision log", "Funding gates"],
    },
}


def build_powerpoint_file(data: dict, output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    guidance = GUIDANCE[data["project_type"]]

    def add_shape(slide, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE):
        item = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        item.fill.solid()
        item.fill.fore_color.rgb = fill
        item.line.fill.background()
        return item

    def add_text(slide, value, x, y, w, h, size, color=INK, bold=False, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = value
        run.font.name = "Aptos Display" if bold else "Aptos"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return box

    def notes(slide, value):
        slide.notes_slide.notes_text_frame.text = value

    def base(kicker, title, number, dark=False):
        slide = prs.slides.add_slide(blank)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = NAVY if dark else WHITE
        add_shape(slide, 0.58, 0.38, 0.18, 0.04, BLUE)
        add_text(slide, kicker.upper(), 0.88, 0.25, 3.0, 0.32, 10, BLUE if not dark else PALE, True)
        add_text(slide, title, 0.58, 0.78, 12.0, 1.05, 30, WHITE if dark else NAVY, True)
        add_shape(slide, 0.58, 7.02, 12.15, 0.01, LINE if not dark else MID)
        add_text(slide, "PROJECT DELIVERY ASSESSMENT", 0.58, 7.08, 4.0, 0.18, 7, MID if not dark else PALE, True)
        add_text(slide, str(number), 12.2, 7.08, 0.5, 0.18, 7, MID if not dark else PALE, True, PP_ALIGN.RIGHT)
        return slide

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_shape(slide, 0, 0, 0.18, 7.5, BLUE)
    add_text(slide, "CLIENT HANDOVER", 0.7, 0.7, 4.0, 0.3, 11, PALE, True)
    add_text(slide, data["title"], 0.7, 1.38, 11.0, 1.2, 38, WHITE, True)
    add_text(slide, "Practical delivery assessment and 30-day control plan.", 0.7, 3.0, 9.0, 0.45, 18, PALE)
    add_shape(slide, 0.7, 4.3, 5.5, 1.4, RGBColor(35, 69, 95))
    add_text(slide, "LIKELY PROJECT TYPE", 0.95, 4.52, 4.5, 0.25, 9, PALE, True)
    add_text(slide, data["project_type"], 0.95, 4.95, 4.8, 0.45, 23, WHITE, True)
    notes(slide, f"Assessment input: {data['assessment']}")

    slide = base("The model", "Four project types need four different control styles", 2)
    types = [
        ("Paint by Numbers", "Outcome known\nMethod known", BLUE, PALE),
        ("Making Movies", "Outcome known\nIntegrated execution", AMBER, RGBColor(246, 240, 228)),
        ("Quest", "Outcome known\nRoute uncertain", GREEN, RGBColor(237, 244, 239)),
        ("Walking in Fog", "Outcome unclear\nRoute unclear", RED, RGBColor(248, 236, 234)),
    ]
    for index, (title, body, accent, fill) in enumerate(types):
        x = 0.58 + index * 3.08
        add_shape(slide, x, 2.32, 2.7, 2.75, fill)
        add_shape(slide, x, 2.32, 2.7, 0.08, accent)
        add_text(slide, title.upper(), x + 0.2, 2.58, 2.3, 0.25, 8, accent, True)
        add_text(slide, title, x + 0.2, 3.18, 2.3, 0.6, 18, NAVY, True)
        add_text(slide, body, x + 0.2, 4.2, 2.3, 0.55, 12, INK)
    add_text(slide, "Match the controls to the uncertainty.", 0.58, 5.65, 11.5, 0.4, 18, BLUE, True)
    notes(slide, "Explain the four categories briefly. The purpose is to select the right controls, not to teach theory.")

    slide = base("The assessment", f"Likely project type: {data['project_type']}", 3)
    add_shape(slide, 0.58, 2.18, 6.45, 3.35, PALE)
    add_text(slide, "PLAIN ENGLISH", 0.9, 2.48, 5.5, 0.25, 9, BLUE, True)
    add_text(slide, guidance["plain"], 0.9, 3.08, 5.6, 1.6, 23, NAVY, True)
    for index, label in enumerate(("Outcome clarity", "Route certainty", "Dependency visibility", "Commitment confidence")):
        y = 2.32 + index * 0.82
        add_shape(slide, 7.75, y, 0.44, 0.44, NAVY, MSO_SHAPE.OVAL)
        add_text(slide, str(index + 1), 7.75, y, 0.44, 0.44, 11, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, label, 8.42, y - 0.03, 3.7, 0.5, 15, NAVY, True)
    notes(slide, f"Use the assessment evidence to explain the call. Input: {data['assessment']}")

    slide = base("Delivery", "What this means for delivery", 4)
    for index, title in enumerate(guidance["delivery"]):
        x = 0.58 + index * 3.08
        add_shape(slide, x, 2.32, 2.7, 2.95, PALE if index == 0 else GREY)
        add_text(slide, str(index + 1).zfill(2), x + 0.2, 2.56, 0.8, 0.35, 11, BLUE, True)
        add_text(slide, title, x + 0.2, 3.35, 2.3, 0.62, 19, NAVY, True)
        add_text(slide, "Use the next review to confirm progress and the next decision.", x + 0.2, 4.38, 2.3, 0.65, 12, INK)
    add_text(slide, "Keep commitments proportionate to confidence.", 0.58, 5.72, 11.5, 0.4, 18, BLUE, True)
    notes(slide, f"For {data['project_type']}, focus delivery on: {', '.join(guidance['delivery'])}.")

    slide = base("Risk", "The main risks are practical and controllable", 5)
    risks = [
        ("Wrong outcome", "The team delivers activity, not the business result.", RED),
        ("Scope locked early", "Commitments get ahead of evidence.", AMBER),
        ("Hidden dependencies", "Data, suppliers or decisions arrive late.", BLUE),
        ("Slow decisions", "Work continues while key choices wait.", NAVY),
        ("False progress", "Tasks close but exposure stays.", GREEN),
    ]
    for index, (title, warning, accent) in enumerate(risks):
        y = 2.05 + index * 0.8
        add_shape(slide, 0.58, y, 0.09, 0.58, accent)
        add_text(slide, title, 0.9, y - 0.02, 3.3, 0.34, 15, NAVY, True)
        add_text(slide, warning, 4.8, y - 0.02, 7.2, 0.34, 13, INK)
        add_shape(slide, 0.9, y + 0.56, 11.1, 0.01, LINE)
    notes(slide, "Focus on risks that change the business result, cost or timing. Escalate against clear triggers.")

    slide = base("Controls", "Use light controls with hard decision points", 6)
    for index, title in enumerate(guidance["control"]):
        x = 0.58 + index * 3.08
        y = 2.32 if index % 2 == 0 else 3.05
        add_shape(slide, x, y, 2.7, 1.9, PALE if index == 2 else GREY)
        add_text(slide, title, x + 0.2, y + 0.35, 2.3, 0.6, 17, NAVY, True, PP_ALIGN.CENTER)
        add_shape(slide, x + 0.95, y + 1.08, 0.8, 0.03, BLUE)
        add_text(slide, "Owner + evidence + date", x + 0.2, y + 1.25, 2.3, 0.35, 11, INK, False, PP_ALIGN.CENTER)
    add_text(slide, "Every control must lead to a decision.", 0.58, 5.65, 11.5, 0.4, 18, BLUE, True)
    notes(slide, f"Recommended controls: {', '.join(guidance['control'])}. Keep reporting short and decision-focused.")

    slide = base("Act now", "Five actions to start this week", 7)
    actions = [
        "Name the accountable sponsor",
        "Write the outcome frame",
        "Rank the top assumptions",
        "Create a two-week evidence plan",
        "Book the first control gate",
    ]
    add_shape(slide, 1.1, 3.12, 10.35, 0.04, LINE)
    for index, action in enumerate(actions):
        x = 0.78 + index * 2.38
        add_shape(slide, x, 2.72, 0.8, 0.8, BLUE if index == 4 else NAVY, MSO_SHAPE.OVAL)
        add_text(slide, str(index + 1), x, 2.72, 0.8, 0.8, 19, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, action, x - 0.58, 4.0, 1.96, 0.9, 13, NAVY, True, PP_ALIGN.CENTER)
    add_text(slide, "Owner + evidence + date + decision", 1.68, 5.55, 9.4, 0.46, 22, BLUE, True, PP_ALIGN.CENTER)
    notes(slide, "Complete or schedule these actions in the first five working days. The gate must end with a decision.")

    slide = base("30 days", "Turn uncertainty into a controllable project", 8, True)
    phases = [
        ("DAYS 1-5", "Frame", "Outcome, sponsor, unknowns"),
        ("DAYS 6-10", "Test", "Users, feasibility, constraints"),
        ("DAYS 11-20", "Shape", "Routes, thin slice, confidence"),
        ("DAYS 21-30", "Commit", "Gate, reclassify, fund next stage"),
    ]
    for index, (days, title, body) in enumerate(phases):
        x = 0.58 + index * 3.08
        fill = BLUE if index == 3 else RGBColor(35, 69, 95)
        add_shape(slide, x, 2.28, 2.7, 2.75, fill)
        add_text(slide, days, x + 0.2, 2.55, 2.3, 0.25, 9, PALE, True)
        add_text(slide, title, x + 0.2, 3.25, 2.3, 0.5, 21, WHITE, True)
        add_text(slide, body, x + 0.2, 4.22, 2.3, 0.55, 12, PALE)
    add_text(slide, "Day 30: choose the delivery model based on evidence.", 0.58, 5.62, 12.0, 0.5, 21, WHITE, True, PP_ALIGN.CENTER)
    notes(slide, "At day 30, reclassify the project and approve the next-stage plan, funding and controls.")

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
